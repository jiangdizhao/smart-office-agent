from __future__ import annotations

import base64
import json
import os
from datetime import UTC, datetime
from email.message import EmailMessage
from pathlib import Path
from typing import Any, Literal

from app.models import ToolResult
from app.presentation_config import REPO_ROOT, presentation_config
from app.tools.presentation_controller import get_presentation_status

GMAIL_COMPOSE_SCOPE = "https://www.googleapis.com/auth/gmail.compose"
SUMMARY_PREFIX = "presentation_summary_"


def _repo_path(value: str) -> Path:
    candidate = Path(value).expanduser()
    if candidate.is_absolute():
        return candidate.resolve()
    return (REPO_ROOT / candidate).resolve()


def gmail_credentials_path() -> Path:
    return _repo_path(
        os.environ.get(
            "SMART_OFFICE_GMAIL_CREDENTIALS",
            "secrets/gmail_credentials.json",
        )
    )


def gmail_token_path() -> Path:
    return _repo_path(
        os.environ.get(
            "SMART_OFFICE_GMAIL_TOKEN",
            "secrets/gmail_token.json",
        )
    )


def _relative_or_absolute(path: Path) -> str:
    try:
        return str(path.relative_to(REPO_ROOT))
    except ValueError:
        return str(path)


def _output_directory() -> Path:
    directory = presentation_config.output_directory.resolve()
    directory.mkdir(parents=True, exist_ok=True)
    return directory


def office_artifact_status() -> dict[str, Any]:
    directory = _output_directory()
    summaries = sorted(
        directory.glob(f"{SUMMARY_PREFIX}*.md"),
        key=lambda path: path.stat().st_mtime,
        reverse=True,
    )
    token_path = gmail_token_path()
    credentials_path = gmail_credentials_path()
    return {
        "output_directory": str(directory),
        "output_directory_relative": _relative_or_absolute(directory),
        "latest_summary": (
            _relative_or_absolute(summaries[0]) if summaries else None
        ),
        "gmail_credentials_path": _relative_or_absolute(credentials_path),
        "gmail_credentials_exists": credentials_path.is_file(),
        "gmail_token_path": _relative_or_absolute(token_path),
        "gmail_token_exists": token_path.is_file(),
        "gmail_draft_configured": credentials_path.is_file() and token_path.is_file(),
        "email_send_enabled": False,
        "recipient_name": presentation_config.recipient_name,
        "recipient_email": presentation_config.recipient_email,
    }


def _slide_text(slide: Any) -> tuple[str, list[str]]:
    title = ""
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None and getattr(title_shape, "has_text_frame", False):
        title = str(title_shape.text).strip()

    blocks: list[str] = []
    for shape in slide.shapes:
        if shape is title_shape or not getattr(shape, "has_text_frame", False):
            continue
        text = " ".join(str(shape.text).split())
        if text:
            blocks.append(text)
    return title, blocks


def _task_lines(task_snapshot: dict[str, Any] | None, language: str) -> list[str]:
    if not task_snapshot:
        return []
    lines: list[str] = []
    for step in task_snapshot.get("steps", []):
        index = step.get("index")
        title = step.get("title") or step.get("tool_name") or "step"
        status = step.get("status") or "unknown"
        if language == "zh":
            lines.append(f"- 第 {index} 步：{title}（{status}）")
        else:
            lines.append(f"- Step {index}: {title} ({status})")
    return lines


def generate_presentation_summary(
    *,
    language: Literal["zh", "en"] = "zh",
    task_snapshot: dict[str, Any] | None = None,
) -> ToolResult:
    source_path = presentation_config.presentation_path.resolve()
    if not source_path.is_file():
        return ToolResult(
            tool_name="office_generate_presentation_summary",
            ok=False,
            message=f"Configured presentation was not found: {source_path}",
            data={
                "execution_mode": "failed",
                "requested_state": {"summary_created": True},
                "presentation_path": str(source_path),
            },
        )

    try:
        from pptx import Presentation

        deck = Presentation(str(source_path))
        slide_records: list[dict[str, Any]] = []
        for index, slide in enumerate(deck.slides, start=1):
            title, blocks = _slide_text(slide)
            slide_records.append(
                {
                    "slide_number": index,
                    "title": title or (f"第 {index} 页" if language == "zh" else f"Slide {index}"),
                    "content": blocks,
                }
            )

        status_result = get_presentation_status()
        presentation_status = dict(status_result.data)
        now = datetime.now(UTC)
        stamp = now.astimezone().strftime("%Y%m%d_%H%M%S_%f")
        output_directory = _output_directory()
        markdown_path = output_directory / f"{SUMMARY_PREFIX}{stamp}.md"
        json_path = output_directory / f"{SUMMARY_PREFIX}{stamp}.json"

        if language == "zh":
            lines = [
                "# 演示摘要",
                "",
                f"- 生成时间：{now.astimezone().isoformat(timespec='seconds')}",
                f"- 演示文件：{source_path.name}",
                f"- 总页数：{len(slide_records)}",
                f"- 当前页：{presentation_status.get('current_slide') or '未放映'}",
                f"- 放映状态：{'正在放映' if presentation_status.get('slideshow_active') else '未放映'}",
                "",
                "## 内容概览",
                "",
            ]
            for record in slide_records:
                lines.append(f"### 第 {record['slide_number']} 页：{record['title']}")
                if record["content"]:
                    for block in record["content"]:
                        lines.append(f"- {block}")
                else:
                    lines.append("- 本页没有可提取的正文文本。")
                lines.append("")
            task_lines = _task_lines(task_snapshot, language)
            if task_lines:
                lines.extend(["## 本次自动化执行", "", *task_lines, ""])
            lines.extend(
                [
                    "## 说明",
                    "",
                    "本摘要由本地 Smart Office Agent 从配置的 PowerPoint 文件和已验证的运行状态生成。",
                    "",
                ]
            )
        else:
            lines = [
                "# Presentation Summary",
                "",
                f"- Generated: {now.astimezone().isoformat(timespec='seconds')}",
                f"- Presentation: {source_path.name}",
                f"- Total slides: {len(slide_records)}",
                f"- Current slide: {presentation_status.get('current_slide') or 'not presenting'}",
                f"- Slide show: {'active' if presentation_status.get('slideshow_active') else 'inactive'}",
                "",
                "## Content overview",
                "",
            ]
            for record in slide_records:
                lines.append(f"### Slide {record['slide_number']}: {record['title']}")
                if record["content"]:
                    for block in record["content"]:
                        lines.append(f"- {block}")
                else:
                    lines.append("- No body text was available for extraction.")
                lines.append("")
            task_lines = _task_lines(task_snapshot, language)
            if task_lines:
                lines.extend(["## Automation run", "", *task_lines, ""])
            lines.extend(
                [
                    "## Note",
                    "",
                    "This summary was generated locally by Smart Office Agent from the configured PowerPoint file and verified runtime state.",
                    "",
                ]
            )

        markdown_text = "\n".join(lines)
        markdown_path.write_text(markdown_text, encoding="utf-8")
        json_payload = {
            "generated_at": now.isoformat(),
            "language": language,
            "presentation_path": str(source_path),
            "presentation_path_relative": _relative_or_absolute(source_path),
            "presentation_status": presentation_status,
            "slides": slide_records,
            "task": task_snapshot,
            "email_send_enabled": False,
        }
        json_path.write_text(
            json.dumps(json_payload, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )

        relative_markdown = _relative_or_absolute(markdown_path)
        relative_json = _relative_or_absolute(json_path)
        return ToolResult(
            tool_name="office_generate_presentation_summary",
            ok=True,
            message=f"Presentation summary created: {relative_markdown}",
            artifacts=[str(markdown_path), str(json_path)],
            data={
                "execution_mode": "real",
                "requested_state": {"summary_created": True},
                "summary_created": True,
                "summary_path": str(markdown_path),
                "summary_path_relative": relative_markdown,
                "summary_json_path": str(json_path),
                "summary_json_path_relative": relative_json,
                "artifact_url": f"/api/office/artifacts/{markdown_path.name}",
                "slide_count": len(slide_records),
                "language": language,
                "presentation_status": presentation_status,
            },
        )
    except Exception as exc:
        return ToolResult(
            tool_name="office_generate_presentation_summary",
            ok=False,
            message=f"Presentation summary could not be generated: {exc}",
            data={
                "execution_mode": "failed",
                "requested_state": {"summary_created": True},
                "error": str(exc),
            },
        )


def latest_summary_path() -> Path | None:
    directory = _output_directory()
    candidates = sorted(
        directory.glob(f"{SUMMARY_PREFIX}*.md"),
        key=lambda path: path.stat().st_mtime,
        reverse=True,
    )
    return candidates[0].resolve() if candidates else None


def _gmail_credentials():
    try:
        from google.auth.transport.requests import Request
        from google.oauth2.credentials import Credentials
    except ImportError as exc:
        raise RuntimeError(
            "Google Gmail dependencies are unavailable. Install backend/requirements-smartoffice.txt."
        ) from exc

    token_path = gmail_token_path()
    credentials_path = gmail_credentials_path()
    if not credentials_path.is_file():
        raise RuntimeError(
            f"Gmail OAuth desktop credentials were not found: {credentials_path}"
        )
    if not token_path.is_file():
        raise RuntimeError(
            "Gmail OAuth token was not found. Run backend/scripts/setup_gmail_oauth.py first."
        )

    credentials = Credentials.from_authorized_user_file(
        str(token_path),
        [GMAIL_COMPOSE_SCOPE],
    )
    if credentials.expired and credentials.refresh_token:
        credentials.refresh(Request())
        token_path.parent.mkdir(parents=True, exist_ok=True)
        token_path.write_text(credentials.to_json(), encoding="utf-8")
    if not credentials.valid:
        raise RuntimeError(
            "Gmail OAuth credentials are invalid. Run backend/scripts/setup_gmail_oauth.py again."
        )
    return credentials


def create_gmail_summary_draft(
    *,
    language: Literal["zh", "en"] = "zh",
    subject: str | None = None,
) -> ToolResult:
    summary_path = latest_summary_path()
    if summary_path is None or not summary_path.is_file():
        return ToolResult(
            tool_name="gmail_create_summary_draft",
            ok=False,
            message="No generated presentation summary is available for the Gmail draft.",
            data={
                "execution_mode": "failed",
                "requested_state": {"gmail_draft_created": True},
                "email_send_enabled": False,
            },
        )

    try:
        from googleapiclient.discovery import build

        credentials = _gmail_credentials()
        summary_text = summary_path.read_text(encoding="utf-8")
        clean_subject = " ".join((subject or "").split())[:180]
        if not clean_subject:
            clean_subject = (
                "演示摘要：Loss.pptx"
                if language == "zh"
                else "Presentation summary: Loss.pptx"
            )

        message = EmailMessage()
        message["To"] = (
            f"{presentation_config.recipient_name} <{presentation_config.recipient_email}>"
        )
        message["Subject"] = clean_subject
        if language == "zh":
            body = (
                f"{presentation_config.recipient_name}，您好：\n\n"
                "以下是 Smart Office Agent 生成的演示摘要。该邮件目前仅保存为 Gmail 草稿，尚未发送。\n\n"
                f"{summary_text}"
            )
        else:
            body = (
                f"Hello {presentation_config.recipient_name},\n\n"
                "Below is the presentation summary generated by Smart Office Agent. "
                "This message has only been saved as a Gmail draft and has not been sent.\n\n"
                f"{summary_text}"
            )
        message.set_content(body)
        raw = base64.urlsafe_b64encode(message.as_bytes()).decode("ascii")

        service = build("gmail", "v1", credentials=credentials, cache_discovery=False)
        draft = (
            service.users()
            .drafts()
            .create(userId="me", body={"message": {"raw": raw}})
            .execute()
        )
        draft_id = str(draft.get("id") or "")
        if not draft_id:
            raise RuntimeError("Gmail API did not return a draft id.")

        return ToolResult(
            tool_name="gmail_create_summary_draft",
            ok=True,
            message=(
                f"Gmail draft created for {presentation_config.recipient_email}; it was not sent."
            ),
            artifacts=[str(summary_path)],
            data={
                "execution_mode": "real",
                "requested_state": {"gmail_draft_created": True},
                "gmail_draft_created": True,
                "gmail_draft_id": draft_id,
                "gmail_drafts_url": "https://mail.google.com/mail/u/0/#drafts",
                "recipient_name": presentation_config.recipient_name,
                "recipient_email": presentation_config.recipient_email,
                "subject": clean_subject,
                "summary_path": str(summary_path),
                "summary_path_relative": _relative_or_absolute(summary_path),
                "email_send_enabled": False,
                "sent": False,
            },
            raw={
                "gmail_draft_id": draft_id,
                "email_send_enabled": False,
                "sent": False,
            },
        )
    except Exception as exc:
        return ToolResult(
            tool_name="gmail_create_summary_draft",
            ok=False,
            message=f"Gmail draft could not be created: {exc}",
            artifacts=[str(summary_path)],
            data={
                "execution_mode": "failed",
                "requested_state": {"gmail_draft_created": True},
                "summary_path": str(summary_path),
                "summary_path_relative": _relative_or_absolute(summary_path),
                "email_send_enabled": False,
                "sent": False,
                "error": str(exc),
            },
        )
